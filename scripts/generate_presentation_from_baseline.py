#!/usr/bin/env python3
"""
Build Real-Time Testing Dashboard briefing deck from the Self-Healing Playwright
technical briefing template (same visual design: fonts, purple/gold palette, layouts).

Requires: python-pptx (see requirements-presentation.txt)
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import sys

_SCRIPTS_DIR = Path(__file__).resolve().parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from generate_presentation_assets import ASSETS_DIR, generate_all  # noqa: E402

BASELINE = Path(__file__).resolve().parents[2].parent / "Self-Healing_Playwright_Framework_Technical_Briefing.pptx"
# Fallback if baseline lives next to project folder
if not BASELINE.exists():
    BASELINE = Path("/Users/manishmohan/Documents/Self-Healing_Playwright_Framework_Technical_Briefing.pptx")

OUT_DIR = Path(__file__).resolve().parent.parent / "presentations"
OUT_FILE = OUT_DIR / "Real-Time_Testing_Dashboard_Technical_Briefing.pptx"

BRAND_FOOTER = "#TheFutureWorksHere"

# Replacement text per slide, in traversal order of non-empty paragraphs (excluding brand footer).
# Paragraph counts must match the baseline deck exactly (see dump in script docstring).
SLIDE_TEXT: list[list[str]] = [
    # Slide 1 (3)
    [
        "REAL-TIME",
        "TESTING DASHBOARD",
        "Live QA observability for CI execution, trends, and Playwright reports",
    ],
    # Slide 2 (12)
    [
        "AGENDA AGENDA",
        "STRATEGIC AGENDA",
        "01",
        "Executive overview & business value",
        "Architecture & data flow",
        "02",
        "Dashboard KPIs & live feed",
        "GitHub Actions ingest",
        "Embedded HTML reports",
        "03",
        "CI pipeline control & deployment",
        "Governance & roadmap",
    ],
    # Slide 3 (11)
    [
        "TECH STACK TECH STACK",
        "MODERN ECOSYSTEM",
        "PLATFORM COMPONENTS",
        "Backend",
        "FastAPI, SQLAlchemy, Alembic, Uvicorn",
        "Frontend",
        "React, TypeScript, Vite",
        "Data",
        "PostgreSQL (Render), SQLite (local)",
        "Hosting",
        "Vercel (UI) + Render (API + DB)",
    ],
    # Slide 4 (11)
    [
        "ARCHITECTURE ARCHITECTURE",
        "SYSTEM DESIGN",
        "LAYERED COMPONENTS",
        "UI presents KPIs, trends, live feed, and embedded",
        "HTML reports.",
        "API LAYER",
        "REST ingest, summary, reports, optional WebSocket,",
        "and GitHub CI trigger.",
        "PERSISTENCE",
        "Postgres stores runs, cases, and optional Playwright",
        "report ZIPs.",
    ],
    # Slide 5 (18)
    [
        "PROCESS FLOW PROCESS",
        "END-TO-END PIPELINE",
        "01",
        "TRIGGER",
        "Start from dashboard (workflow_dispatch) or ",
        "push to GitHub.",
        "02",
        "CI EXECUTION",
        "Playwright suite runs on GitHub Actions; jobs ",
        "and steps visible in UI.",
        "03",
        "INGEST",
        "Workflow POSTs JSON results and optional ",
        "report_zip to the API.",
        "04",
        "OBSERVE",
        "Dashboard updates KPIs, trends, feed, and ",
        "embedded HTML viewer.",
    ],
    # Slide 6 (21)
    [
        "DASHBOARD DASHBOARD DASHBOARD",
        "LIVE KPI SURFACE",
        "EXECUTIVE VISIBILITY",
        "SINGLE PANE OF GLASS",
        "Totals for runs, cases, pass rate, and open",
        "defects in one view.",
        "Trends refresh as new CI runs are ingested.",
        "Portfolio roll-ups without opening GitHub Actions.",
        "NO LOG ARCHAEOLOGY",
        "Leadership and QA stop chasing artifacts across",
        "CI tabs and storage buckets.",
        "ALWAYS CURRENT",
        "Latest CI runs surface first when",
        "DATA_SOURCE=github.",
        "DEMO READY",
        "Synthetic runs available in demo mode for",
        "stakeholder walkthroughs.",
        "KPI EXAMPLES",
        "Total Runs · Pass Rate · Module quality",
        "Example:",
        "https://realtime-testing-dashboard-ooka.vercel.app",
    ],
    # Slide 7 (23)
    [
        "INGEST INGEST",
        "GITHUB ACTIONS PUBLISH",
        "SECURE CI PATH (PURPLE)",
        "STRUCTURED PAYLOAD",
        "suite_name, environment, build_version,",
        "test_cases[] from Playwright JSON.",
        "HTML REPORT ZIP",
        "Zip playwright-report/ and POST multipart",
        "run-with-report.",
        "SERVER-ONLY SECRET",
        "GITHUB_ACTIONS_INGEST_TOKEN never exposed to",
        "the browser.",
        "AUTOMATIC REFRESH",
        "Summary and live feed update after",
        "successful ingest.",
        "ENDPOINTS",
        "POST /api/ingest/github-actions/run",
        "POST /api/ingest/github-actions/run-with-report",
        "Triggers: After every CI run for the dashboard.",
        "Efficiency: Non-blocking POST from workflow.",
        "Accuracy: Same JSON schema as Playwright output.",
        "Repo: mmohansqaai/SelfHealingPlaywrightFramework",
        "Pair with examples/github-actions-publish-step.yml",
    ],
    # Slide 8 (9)
    [
        "HTML REPORT HTML REPORT",
        "EMBEDDED VIEWER",
        "IN-DASHBOARD PLAYWRIGHT REPORTS",
        "ZIP Storage",
        "API stores report ZIP and serves /api/runs/{id}/report/... so teams browse without a public artifact URL.",
        "Inline View",
        "Managers and leads open the dashboard; engineers can still open full Playwright HTML for forensics.",
        "Theme Alignment",
        "Report shell styled to match dashboard dark theme when served from the API.",
    ],
    # Slide 9 (24)
    [
        "CONFIGURATION CONFIGURATION",
        "OPERATING MODES",
        "RUNTIME ORDER",
        "demo",
        "↓",
        "github",
        "↓",
        "INGEST",
        "ENVIRONMENT",
        "Global:",
        "Set",
        "DATA_SOURCE",
        "env variable (demo vs github).",
        "Per-Deploy:",
        "Override DATABASE_URL and CORS_ORIGINS on",
        "Render for each Vercel origin and API host.",
        "Keep GITHUB_ACTIONS_INGEST_TOKEN server-side only.",
        "DEDUPLICATION",
        "Scoring:",
        "GITHUB_CI_TOKEN enables workflow_dispatch",
        "from the dashboard.",
        "Uniqueness:",
        "Vercel rewrites /api to Render so the",
        "browser stays same-origin.",
    ],
    # Slide 10 (11)
    [
        "PIPELINE PIPELINE",
        "CI/CD & DASHBOARD",
        "AUTOMATED LOOP",
        "Dashboard triggers GitHub Actions; Playwright runs on",
        "every workflow_dispatch or push.",
        "EXECUTION VISIBILITY",
        "Poll GitHub for job/step status (queued, in_progress,",
        "completed) in the CI panel.",
        "DASHBOARD PUBLISH",
        "Workflow posts JSON (+ optional ZIP) to Render API;",
        "UI refreshes summary and live feed.",
    ],
    # Slide 11 (17)
    [
        "OBSERVABILITY OBSERVABILITY",
        "DASHBOARD VS. FORENSIC REPORTS",
        "FEATURE",
        "REAL-TIME DASHBOARD",
        "PLAYWRIGHT HTML REPORT",
        "Audience",
        "Leadership / QA Managers",
        "Engineers / SDETs",
        "Focus",
        "Trends, Pass/Fail, Build History, CI control",
        "Traces, Videos, Screenshots",
        "Report Detail",
        "Run list, KPIs, embedded HTML when uploaded",
        "All attempts, step timings, attachments",
        "Actionable Info",
        "Portfolio health & pipeline status",
        "Forensic debugging of individual failures",
    ],
    # Slide 12 (15)
    [
        "GOVERNANCE ROADMAP",
        "GOVERNANCE & ROADMAP",
        "CORE GOVERNANCE",
        "Secrets on server:",
        "Ingest token, GitHub PAT, and database URL never ship to the browser.",
        "CORS allowlist:",
        "Each Vercel origin must be listed on Render for production.",
        "FUTURE ROADMAP",
        "Alerting:",
        "Slack/email when pass rate drops below thresholds.",
        "Multi-repo:",
        "Support several GitHub repos and RBAC per team.",
        "STRATEGIC VISION",
        "Single quality cockpit for Playwright CI",
        "across the portfolio.",
    ],
    # Slide 13 (3)
    [
        "THANK",
        "YOU",
        "QUESTIONS?",
    ],
]


def _set_paragraph_text(paragraph, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def _paragraphs_to_replace(slide) -> list:
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            current = "".join(run.text for run in paragraph.runs).strip()
            if not current or current == BRAND_FOOTER:
                continue
            out.append(paragraph)
    return out


# (slide_number 1-based or None for all slides, shape name substring, asset filename)
IMAGE_RULES: list[tuple[int | None, str, str]] = [
    (None, "9.26.42", "footer.png"),
    (None, "Google Shape;129", "logo.png"),
    (None, "Image 0", "slide_bg.png"),
    (1, "Image 1", "hero_dashboard.png"),
    (3, "Image 1", "tech_ecosystem.png"),
    (4, "SelfHealing", "architecture.png"),
    (4, "Architecture", "architecture.png"),
    (5, "6.34.21", "process_flow.png"),
    (7, "7.46.09", "ingest_diagram.png"),
    (8, "7.46.53", "html_report_embed.png"),
    (10, "7.52.04", "ci_steps.png"),
    (10, "8.22.24", "github_workflow.png"),
    (11, "8.23.29", "dashboard_comparison.png"),
    (13, "9.22.49", "thank_you_hero.png"),
]


def _replace_picture_blob(shape, image_path: Path) -> None:
    """Swap embedded image bytes while keeping position and size."""
    blip = shape._element.blipFill.blip
    r_id = blip.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    )
    image_part = shape.part.related_parts[r_id]
    image_part.blob = image_path.read_bytes()


def apply_slide_images(slide, slide_number: int, assets: dict[str, Path]) -> int:
    replaced = 0
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        name = shape.name or ""
        for rule_slide, name_part, asset_name in IMAGE_RULES:
            if rule_slide is not None and rule_slide != slide_number:
                continue
            if name_part not in name:
                continue
            asset_path = assets.get(asset_name) or ASSETS_DIR / asset_name
            if not asset_path.exists():
                raise FileNotFoundError(f"Missing asset: {asset_path}")
            _replace_picture_blob(shape, asset_path)
            replaced += 1
            break
    return replaced


def apply_slide_content(slide, replacements: list[str]) -> None:
    paragraphs = _paragraphs_to_replace(slide)
    if len(paragraphs) != len(replacements):
        raise ValueError(
            f"Slide {slide.slide_id}: expected {len(replacements)} text blocks, found {len(paragraphs)}. "
            "Update SLIDE_TEXT in generate_presentation_from_baseline.py."
        )
    for paragraph, new_text in zip(paragraphs, replacements):
        _set_paragraph_text(paragraph, new_text)


def build() -> Path:
    if not BASELINE.exists():
        raise FileNotFoundError(f"Baseline deck not found: {BASELINE}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(BASELINE, OUT_FILE)

    assets = generate_all()
    prs = Presentation(str(OUT_FILE))
    if len(prs.slides) != len(SLIDE_TEXT):
        raise ValueError(f"Baseline has {len(prs.slides)} slides; content map has {len(SLIDE_TEXT)}")

    total_images = 0
    for index, (slide, replacements) in enumerate(zip(prs.slides, SLIDE_TEXT), start=1):
        apply_slide_content(slide, replacements)
        total_images += apply_slide_images(slide, index, assets)

    prs.core_properties.title = "Real-Time Testing Dashboard — Technical Briefing"
    prs.core_properties.subject = "QA observability dashboard for live CI and Playwright"
    prs.save(str(OUT_FILE))
    print(f"Replaced {total_images} picture(s) with dashboard assets from {ASSETS_DIR}")
    return OUT_FILE


def main() -> None:
    path = build()
    print(f"Wrote {path}")


if __name__ == "__main__":
    main()
